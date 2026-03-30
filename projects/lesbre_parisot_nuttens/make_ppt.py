"""Generate presentation.pptx for the GDELT × Brent project."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── colours ──────────────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x1A, 0x2D, 0x4A)   # slide backgrounds / title bar
MID_BLUE  = RGBColor(0x2E, 0x5E, 0x9B)   # accents
ORANGE    = RGBColor(0xE8, 0x7A, 0x22)   # highlights
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ── helpers ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


def bg(slide, colour: RGBColor = DARK_BLUE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_textbox(slide, text: str, left, top, width, height,
                font_size=18, bold=False, colour=WHITE,
                align=PP_ALIGN.LEFT, wrap=True) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = colour


def add_rect(slide, left, top, width, height, fill_colour: RGBColor,
             line_colour: RGBColor | None = None):
    shape = slide.shapes.add_shape(
        1, left, top, width, height)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if line_colour:
        shape.line.color.rgb = line_colour
    else:
        shape.line.fill.background()
    return shape


def set_notes(slide, text: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def header_bar(slide, title: str, subtitle: str = ""):
    """Dark bar at top with slide title."""
    add_rect(slide, 0, 0, W, Inches(1.15), DARK_BLUE)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.18), Inches(11), Inches(0.65),
                font_size=28, bold=True, colour=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.78), Inches(11), Inches(0.35),
                    font_size=14, colour=ORANGE)


def bullet_block(slide, items: list[str], left, top, width, height,
                 font_size=17, colour=DARK_GREY, bullet="▸ "):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = bullet + item
        run.font.size  = Pt(font_size)
        run.font.color.rgb = colour


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Title."""
    s = blank_slide(prs)
    bg(s, DARK_BLUE)

    # accent stripe
    add_rect(s, 0, Inches(3.0), Inches(0.12), Inches(1.9), ORANGE)

    add_textbox(s, "GDELT × Brent",
                Inches(0.35), Inches(1.5), Inches(12), Inches(1.1),
                font_size=54, bold=True, colour=WHITE, align=PP_ALIGN.LEFT)

    add_textbox(s, "Geopolitical Tension as a Predictive Signal for Oil Prices",
                Inches(0.35), Inches(2.55), Inches(10), Inches(0.7),
                font_size=22, colour=ORANGE, align=PP_ALIGN.LEFT)

    add_textbox(s,
                "Lesbre · Parisot · Nuttens  |  ML for Finance — Alternative Data",
                Inches(0.35), Inches(3.35), Inches(10), Inches(0.45),
                font_size=15, colour=RGBColor(0xAA, 0xBB, 0xCC))

    add_textbox(s, "2015 – 2024  |  Binary classification  |  J+3 horizon",
                Inches(0.35), Inches(3.85), Inches(10), Inches(0.4),
                font_size=13, colour=RGBColor(0x88, 0x99, 0xAA))

    set_notes(s, """\
[~30 s]
Bonjour à tous. Notre projet s'appelle GDELT × Brent.
La question centrale : est-ce qu'un score de tension géopolitique construit à partir de données \
médiatiques peut prédire la direction des prix du pétrole Brent à 3 jours ?
On va vous montrer comment on a construit ce signal, ce qu'on a trouvé — et surtout, \
ce qu'on n'a pas trouvé, parce que c'est là que c'est intéressant.
""")


def slide_question(prs):
    """Slide 2 — Research question & intuition."""
    s = blank_slide(prs)
    bg(s, LIGHT_GREY)
    header_bar(s, "The Question", "Why would geopolitics move oil prices?")

    # Left panel — hypothesis
    add_rect(s, Inches(0.3), Inches(1.3), Inches(5.8), Inches(5.7), WHITE)
    add_textbox(s, "Hypothesis",
                Inches(0.5), Inches(1.45), Inches(5.4), Inches(0.5),
                font_size=17, bold=True, colour=MID_BLUE)
    bullet_block(s, [
        "Middle East supplies ~35% of global oil",
        "Geopolitical shocks historically spike Brent",
        "GDELT encodes media events in real-time → alternative data",
        "If the signal exists, can ML find it before markets price it in?",
    ], Inches(0.45), Inches(2.0), Inches(5.6), Inches(4.5),
        font_size=16, colour=DARK_GREY)

    # Right panel — task definition
    add_rect(s, Inches(7.2), Inches(1.3), Inches(5.8), Inches(5.7), WHITE)
    add_textbox(s, "ML Task",
                Inches(7.4), Inches(1.45), Inches(5.4), Inches(0.5),
                font_size=17, bold=True, colour=MID_BLUE)
    bullet_block(s, [
        "Binary classification",
        "Target: will Brent close higher in 3 trading days?",
        "Features: daily GDELT aggregates over the Middle East",
        "Period: 2015 – 2024  (2 508 trading days)",
        "Not a regression — we predict direction, not magnitude",
    ], Inches(7.35), Inches(2.0), Inches(5.6), Inches(4.5),
        font_size=16, colour=DARK_GREY)

    # divider label
    add_textbox(s, "vs", Inches(6.05), Inches(3.8), Inches(1.1), Inches(0.6),
                font_size=26, bold=True, colour=ORANGE, align=PP_ALIGN.CENTER)

    set_notes(s, """\
[~50 s]
Le Brent est le benchmark mondial du pétrole. Le Moyen-Orient — Iran, Irak, Arabie Saoudite, \
Yémen, Syrie, Israël, Palestine, Émirats, Koweït — représente environ 35% de la production \
mondiale et une part encore plus grande des exportations maritimes.

L'idée : les chocs géopolitiques dans cette région ont historiquement provoqué des mouvements \
brusques sur le Brent. La question c'est : est-ce que cette relation est détectable dans les \
données médiatiques AVANT que les marchés ne l'intègrent complètement ?

On formule ça comme une classification binaire : est-ce que le Brent va clôturer plus haut \
dans 3 jours trading ? Pas une régression sur le niveau des prix — juste la direction.
""")


def slide_data(prs):
    """Slide 3 — Data sources."""
    s = blank_slide(prs)
    bg(s, LIGHT_GREY)
    header_bar(s, "Data", "Two public sources — no proprietary data")

    # GDELT box
    add_rect(s, Inches(0.3), Inches(1.3), Inches(6.0), Inches(5.7), WHITE)
    add_rect(s, Inches(0.3), Inches(1.3), Inches(6.0), Inches(0.5),
             MID_BLUE)
    add_textbox(s, "GDELT 1.0 — Global Event Database",
                Inches(0.4), Inches(1.33), Inches(5.8), Inches(0.45),
                font_size=14, bold=True, colour=WHITE)
    bullet_block(s, [
        "Media events worldwide → CAMEO taxonomy",
        "Public CSV files, updated daily",
        "3 653 daily files downloaded (2015 – 2024)",
        "Filter: 9 ME country codes (IRN, IRQ, SAU, ISR…)",
        "Key signals: Goldstein Scale (−10→+10), AvgTone",
        "~12 600 ME events per trading day on average",
    ], Inches(0.45), Inches(1.9), Inches(5.7), Inches(4.8),
        font_size=15, colour=DARK_GREY)

    # Brent box
    add_rect(s, Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.7), WHITE)
    add_rect(s, Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
             MID_BLUE)
    add_textbox(s, "Brent Crude — BZ=F via yfinance",
                Inches(7.1), Inches(1.33), Inches(5.8), Inches(0.45),
                font_size=14, bold=True, colour=WHITE)
    bullet_block(s, [
        "Continuous futures, daily OHLCV",
        "2 508 trading days (2015-01-05 → 2024-12-24)",
        "Target: Close(t+3) > Close(t)",
        "Class balance: 53.2% Up / 46.8% Down",
        "Slight upward drift → mild imbalance expected",
    ], Inches(7.15), Inches(1.9), Inches(5.7), Inches(4.8),
        font_size=15, colour=DARK_GREY)

    set_notes(s, """\
[~45 s]
Deux sources de données entièrement publiques.

GDELT : on a téléchargé 3 653 fichiers CSV quotidiens, filtrés sur 9 codes pays du \
Moyen-Orient. Chaque ligne est un événement encodé selon la taxonomie CAMEO. \
On utilise principalement la Goldstein Scale — qui mesure l'impact déstabilisant théorique \
d'un événement entre −10 et +10 — et l'AvgTone, le sentiment moyen de la couverture médiatique.

Brent : données Yahoo Finance via yfinance. La cible est binaire : est-ce que le prix \
de clôture dans 3 jours trading est supérieur au prix d'aujourd'hui ? \
Les classes sont légèrement déséquilibrées — 53% hausse contre 47% baisse — \
ce qui reflète la tendance haussière du Brent sur la période.
""")


def slide_features(prs):
    """Slide 4 — Feature engineering."""
    s = blank_slide(prs)
    bg(s, LIGHT_GREY)
    header_bar(s, "Feature Engineering", "12 daily aggregates from GDELT")

    # table header
    for col, x, w in [("Feature", 0.3, 3.5), ("Construction", 3.9, 4.5), ("Intuition", 8.5, 4.6)]:
        add_rect(s, Inches(x), Inches(1.3), Inches(w), Inches(0.38), MID_BLUE)
        add_textbox(s, col, Inches(x + 0.1), Inches(1.32), Inches(w - 0.1), Inches(0.36),
                    font_size=13, bold=True, colour=WHITE)

    rows = [
        ("n_events",          "COUNT(*) ME events",                 "Total geopolitical activity"),
        ("n_conflict_events", "COUNT WHERE code ~ 18x/19x/20x",     "Direct conflict intensity"),
        ("conflict_ratio",    "n_conflict / n_events",               "Proportion of hostile events"),
        ("goldstein_mean",    "MEAN(GoldsteinScale)",                "Average stability signal"),
        ("avg_tone",          "MEAN(AvgTone)",                       "Media framing (hostile < 0)"),
        ("n_mentions",        "SUM(NumMentions)",                    "Media intensity / attention"),
        ("n_articles",        "SUM(NumArticles)",                    "Breadth of coverage"),
        ("goldstein_7d_ma",   "7-day rolling mean of goldstein",     "Smoothed trend"),
        ("tension_spike",     "goldstein < 7d_ma − 1.5σ",           "Binary escalation flag"),
        ("mentions_zscore",   "(n_mentions − 7d_ma) / σ",           "Abnormal media surge"),
    ]

    stripe = [WHITE, LIGHT_GREY]
    for i, (feat, constr, intuit) in enumerate(rows):
        y = Inches(1.68 + i * 0.48)
        h = Inches(0.48)
        for x, w in [(0.3, 3.5), (3.9, 4.5), (8.5, 4.6)]:
            add_rect(s, Inches(x), y, Inches(w), h, stripe[i % 2])
        add_textbox(s, feat,    Inches(0.4),  y + Inches(0.08), Inches(3.3),  h,
                    font_size=12, bold=True,  colour=MID_BLUE)
        add_textbox(s, constr,  Inches(4.0),  y + Inches(0.08), Inches(4.3),  h,
                    font_size=11, colour=DARK_GREY)
        add_textbox(s, intuit,  Inches(8.6),  y + Inches(0.08), Inches(4.4),  h,
                    font_size=11, colour=DARK_GREY)

    add_textbox(s, "Note: goldstein_min = −10.0 on every day → excluded (constant, no variance)",
                Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.35),
                font_size=11, colour=RGBColor(0x88, 0x44, 0x00))

    set_notes(s, """\
[~45 s]
On agrège les événements GDELT par jour calendar pour construire 10 features. \
Il y a deux familles :

Les features de volume — n_events, n_mentions, n_articles — qui mesurent l'intensité \
de l'activité médiatique. Elles sont très corrélées entre elles (r = 0.93).

Les features de type — conflict_ratio, goldstein_mean, avg_tone — qui mesurent la nature \
des événements, pas leur volume. conflict_ratio et goldstein_mean sont fortement corrélés \
négativement (r = −0.91) ce qui est attendu : plus d'événements conflictuels = Goldstein plus bas.

Point important : goldstein_min est constant à −10.0 sur toute la période. \
Chaque jour, au moins un événement maximalement déstabilisant est reporté au Moyen-Orient. \
Cette feature n'apporte aucune information et a été exclue.
""")


def slide_method(prs):
    """Slide 5 — Methodology."""
    s = blank_slide(prs)
    bg(s, LIGHT_GREY)
    header_bar(s, "Methodology", "The single most important design choice")

    # big highlight box
    add_rect(s, Inches(0.3), Inches(1.35), Inches(12.7), Inches(1.3),
             RGBColor(0xFF, 0xF3, 0xE0))
    add_rect(s, Inches(0.3), Inches(1.35), Inches(0.1), Inches(1.3), ORANGE)
    add_textbox(s,
                "TimeSeriesSplit — never KFold",
                Inches(0.55), Inches(1.4), Inches(12), Inches(0.5),
                font_size=20, bold=True, colour=DARK_BLUE)
    add_textbox(s,
                "KFold shuffles observations → model trains on data posterior to the test set → look-ahead bias → "
                "inflated accuracy. TimeSeriesSplit always trains on the past, evaluates on the future.",
                Inches(0.55), Inches(1.88), Inches(12.5), Inches(0.7),
                font_size=13, colour=DARK_GREY)

    # two model boxes
    for x, title, desc in [
        (0.3, "Logistic Regression",
         "Interpretable baseline.\nCoefficients show which features drive the prediction.\nFast, no hyperparameters to tune."),
        (6.8, "Random Forest",
         "Captures non-linearities.\nFeature importance out of the box.\n100 trees, random_state=42."),
    ]:
        add_rect(s, Inches(x), Inches(2.85), Inches(6.2), Inches(2.4), WHITE)
        add_rect(s, Inches(x), Inches(2.85), Inches(6.2), Inches(0.42), MID_BLUE)
        add_textbox(s, title, Inches(x + 0.15), Inches(2.88),
                    Inches(5.9), Inches(0.38), font_size=15, bold=True, colour=WHITE)
        add_textbox(s, desc, Inches(x + 0.15), Inches(3.3),
                    Inches(5.9), Inches(1.8), font_size=14, colour=DARK_GREY)

    # benchmarks
    add_rect(s, Inches(0.3), Inches(5.45), Inches(12.7), Inches(1.7), WHITE)
    add_textbox(s, "Benchmarks to beat",
                Inches(0.5), Inches(5.5), Inches(12), Inches(0.4),
                font_size=15, bold=True, colour=MID_BLUE)
    bullet_block(s, [
        "Majority class — always predict 'Up' (most frequent) → 53.9% accuracy on test set",
        "Rolling volatility — go long when recent vol < threshold → 48.1% accuracy on test set",
    ], Inches(0.45), Inches(5.95), Inches(12.5), Inches(1.1), font_size=14, colour=DARK_GREY)

    set_notes(s, """\
[~50 s]
Le choix méthodologique le plus important : TimeSeriesSplit, pas KFold standard.

KFold mélange les observations aléatoirement. Sur des séries temporelles, ça permet \
au modèle de s'entraîner sur des données qui sont chronologiquement postérieures \
au set de test — c'est un biais de look-ahead direct. Les résultats seraient gonflés \
et sans valeur. TimeSeriesSplit garantit qu'on s'entraîne toujours sur le passé \
et qu'on évalue sur le futur. 5 folds.

Deux modèles : Logistic Regression pour l'interprétabilité, Random Forest pour \
les non-linéarités et les importances de features.

Deux benchmarks naïfs à battre : le classifieur majoritaire — qui prédit toujours \
"hausse" — atteint 53.9% de précision sur le test set. Le signal de volatilité rolling \
fait 48.1%. Un modèle qui ne bat pas ces benchmarks n'apporte rien.
""")


def slide_results(prs):
    """Slide 6 — Results."""
    s = blank_slide(prs)
    bg(s, LIGHT_GREY)
    header_bar(s, "Results", "What the numbers say — and don't say")

    # accuracy table
    add_rect(s, Inches(0.3), Inches(1.35), Inches(6.0), Inches(5.7), WHITE)
    add_textbox(s, "CV Accuracy & F1 (macro)",
                Inches(0.5), Inches(1.42), Inches(5.6), Inches(0.4),
                font_size=15, bold=True, colour=MID_BLUE)

    table_rows = [
        ("Model",                "Accuracy",  "F1 macro", True),
        ("Majority Class ★",     "54.6%",     "0.35",     False),
        ("Logistic Regression",  "52.7%",     "0.45",     False),
        ("Random Forest",        "50.8%",     "0.49",     False),
        ("Rolling Volatility",   "47.2%",     "0.47",     False),
    ]
    for i, (name, acc, f1, is_hdr) in enumerate(table_rows):
        y = Inches(1.85 + i * 0.55)
        colour = MID_BLUE if is_hdr else (RGBColor(0xE8, 0xF4, 0xFD) if i % 2 == 0 else WHITE)
        add_rect(s, Inches(0.3), y, Inches(6.0), Inches(0.54), colour)
        fc = WHITE if is_hdr else DARK_GREY
        add_textbox(s, name, Inches(0.45), y + Inches(0.09), Inches(3.0), Inches(0.4),
                    font_size=13, bold=is_hdr, colour=fc)
        add_textbox(s, acc,  Inches(3.6),  y + Inches(0.09), Inches(1.2), Inches(0.4),
                    font_size=13, bold=is_hdr, colour=fc, align=PP_ALIGN.CENTER)
        add_textbox(s, f1,   Inches(4.9),  y + Inches(0.09), Inches(1.2), Inches(0.4),
                    font_size=13, bold=is_hdr, colour=fc, align=PP_ALIGN.CENTER)

    # key take-aways
    add_rect(s, Inches(6.5), Inches(1.35), Inches(6.6), Inches(5.7), WHITE)
    add_textbox(s, "Key takeaways",
                Inches(6.7), Inches(1.42), Inches(6.2), Inches(0.4),
                font_size=15, bold=True, colour=MID_BLUE)
    bullet_block(s, [
        "Neither model beats majority class on accuracy",
        "But both beat it on F1 macro — they predict both classes, not just 'Up'",
        "Accuracy spread across CV folds is narrow → weak but consistent signal, not a one-crisis artefact",
        "Backtest (LR, 2023-24): −22.8% vs buy-and-hold −12.8%",
        "Honest conclusion: no tradeable edge at J+3 in current form",
    ], Inches(6.65), Inches(1.95), Inches(6.3), Inches(5.0), font_size=15, colour=DARK_GREY)

    set_notes(s, """\
[~55 s]
Les résultats sont clairs — et honnêtes.

En accuracy sur validation croisée, aucun modèle ne bat le classifieur majoritaire. \
LR fait 52.7% contre 54.6% pour le majority class. Random Forest fait 50.8%.

Mais regardez le F1 macro. Le majority class a un F1 de 0.35 — parce qu'il ne prédit \
jamais la classe minoritaire, donc son recall sur les baisses est zéro. \
LR fait 0.45 et RF fait 0.49 — ils reconnaissent les deux directions.

La distribution des précisions par fold est étroite — entre 49 et 56% pour LR. \
Ce n'est pas un artefact concentré sur une crise particulière : le signal est faible \
mais consistant sur toute la période.

Le backtest confirme : le signal LR sous-performe le buy-and-hold sur 2023-2024. \
Conclusion honnête : on n'a pas de signal tradeable au J+3 dans cette configuration.
""")


def slide_finding(prs):
    """Slide 7 — The most interesting finding."""
    s = blank_slide(prs)
    bg(s, LIGHT_GREY)
    header_bar(s, "The Most Interesting Finding",
               "It's not that the model fails — it's why")

    add_rect(s, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.7), WHITE)

    add_textbox(s, "Feature importance is uniform — and that tells us something",
                Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
                font_size=18, bold=True, colour=MID_BLUE)

    # feature importance approximation (text-based)
    feats = [
        ("goldstein_7d_ma",  11.5, "Sustained deterioration trend"),
        ("avg_tone",         10.9, "Hostile media framing"),
        ("mentions_zscore",  10.4, "Abnormal media surge"),
        ("conflict_ratio",   10.2, "Proportion of hostile events"),
        ("goldstein_mean",    9.97, "Daily stability signal"),
        ("n_conflict_events", 9.83, "Raw conflict event count"),
    ]
    for i, (feat, imp, desc) in enumerate(feats):
        y = Inches(2.1 + i * 0.62)
        bar_w = Inches(imp * 0.22)
        add_rect(s, Inches(0.5), y, bar_w, Inches(0.38), MID_BLUE)
        add_textbox(s, f"{feat}  ({imp:.1f}%)",
                    Inches(0.5), y, Inches(4.5), Inches(0.38),
                    font_size=12, bold=True, colour=WHITE)
        add_textbox(s, desc, Inches(5.1), y, Inches(7.8), Inches(0.38),
                    font_size=12, colour=DARK_GREY)

    add_rect(s, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.9),
             RGBColor(0xFF, 0xF3, 0xE0))
    add_rect(s, Inches(0.3), Inches(6.0), Inches(0.08), Inches(0.9), ORANGE)
    add_textbox(s,
                "Markets likely price in GDELT signals fast. The J+3 horizon may capture "
                "mean-reversion after overreaction, not a genuine predictive lag. "
                "This is precisely what an efficient market would predict.",
                Inches(0.55), Inches(6.05), Inches(12.4), Inches(0.8),
                font_size=13, colour=DARK_GREY)

    set_notes(s, """\
[~55 s]
Le résultat le plus intéressant, ce n'est pas que les modèles ne battent pas le benchmark. \
C'est pourquoi ils ne le battent pas.

Les importances de features du Random Forest sont remarquablement uniformes — \
entre 9.8% et 11.5% pour les 6 features principales. Aucune feature ne domine. \
Ce n'est pas un modèle qui a trouvé une relation claire : c'est un modèle qui distribue \
son attention sur des signaux tous également faibles.

Les trois features les plus importantes ont une logique claire : \
goldstein_7d_ma capture une détérioration *soutenue* et non le bruit quotidien ; \
avg_tone reflète le cadrage hostile de la couverture ; \
mentions_zscore détecte une sur-attention médiatique anormale. \
Ce sont des signaux de *changement de régime*, pas d'événements ponctuels.

L'explication probable : GDELT est public et gratuit. Si des traders quantitatifs \
l'utilisent déjà — et c'est vraisemblable — l'edge est arbitragé. \
Le J+3 ne capture peut-être pas une anticipation mais une mean-reversion après surréaction. \
C'est exactement ce que prédirait un marché efficient.
""")


def slide_limitations(prs):
    """Slide 8 — Limitations."""
    s = blank_slide(prs)
    bg(s, LIGHT_GREY)
    header_bar(s, "Limitations", "What this project cannot claim")

    items = [
        ("Look-ahead bias",
         "Protected by TimeSeriesSplit + shift(-3) target. No training on future data."),
        ("GDELT = media, not events",
         "An unreported conflict generates no signal. A minor but viral event generates a strong one."),
        ("Coverage non-stationarity",
         "GDELT indexes far more articles in 2023 than 2015. Raw volume features are not comparable "
         "across time. Partially mitigated by conflict_ratio and mentions_zscore."),
        ("Frictionless backtest",
         "No transaction costs, no bid-ask spread, no position sizing. Returns shown are an upper bound."),
        ("Structural breaks",
         "COVID 2020 and Ukraine 2022 moved Brent for reasons orthogonal to Middle East geopolitics. "
         "Including these periods may hurt rather than help generalisation."),
    ]

    for i, (title, desc) in enumerate(items):
        y = Inches(1.4 + i * 1.12)
        add_rect(s, Inches(0.3), y, Inches(12.7), Inches(1.05), WHITE if i % 2 == 0 else LIGHT_GREY)
        add_rect(s, Inches(0.3), y, Inches(0.08), Inches(1.05), ORANGE)
        add_textbox(s, title, Inches(0.55), y + Inches(0.06),
                    Inches(3.2), Inches(0.42), font_size=14, bold=True, colour=MID_BLUE)
        add_textbox(s, desc,  Inches(3.9), y + Inches(0.06),
                    Inches(9.0), Inches(0.9), font_size=13, colour=DARK_GREY)

    set_notes(s, """\
[~40 s]
Quelques limitations importantes à reconnaître explicitement.

Le biais de look-ahead : on l'a protégé avec TimeSeriesSplit et le shift(-3). Mais on en parle \
parce que c'est le risque numéro un dans ce type de projet, et montrer qu'on y a pensé est un point positif.

GDELT encode la couverture médiatique, pas les événements réels. Un conflit non-reporté \
génère un signal nul même si son impact marché est large.

La non-stationnarité de couverture : GDELT est beaucoup plus dense en 2023 qu'en 2015. \
On ne peut pas comparer des n_mentions bruts sur toute la période. Les features normalisées \
comme conflict_ratio et mentions_zscore atténuent le problème.

Le backtest est frictionless — c'est un upper bound, pas un P&L réel.

Et les ruptures structurelles : COVID 2020 et Ukraine 2022 ont fait bouger le Brent \
pour des raisons sans lien avec la géopolitique moyen-orientale.
""")


def slide_conclusion(prs):
    """Slide 9 — Conclusion & next steps."""
    s = blank_slide(prs)
    bg(s, DARK_BLUE)

    add_textbox(s, "Conclusion",
                Inches(0.4), Inches(0.8), Inches(12), Inches(0.7),
                font_size=34, bold=True, colour=WHITE)

    add_rect(s, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.06), ORANGE)

    # two column summary
    add_textbox(s, "What we found",
                Inches(0.4), Inches(1.75), Inches(6.0), Inches(0.5),
                font_size=17, bold=True, colour=ORANGE)
    bullet_block(s, [
        "GDELT geopolitical tension is measurable and visually correlated with Brent spikes",
        "Feature importance is diffuse — no single signal dominates",
        "Models outperform on F1 but not on accuracy vs majority class",
        "No tradeable edge at J+3 in current form",
    ], Inches(0.4), Inches(2.3), Inches(6.1), Inches(3.5),
        font_size=15, colour=WHITE, bullet="✓ ")

    add_textbox(s, "What would come next",
                Inches(7.0), Inches(1.75), Inches(6.0), Inches(0.5),
                font_size=17, bold=True, colour=ORANGE)
    bullet_block(s, [
        "Test J+1 vs J+3 vs J+5 to characterise signal decay",
        "Actor-level disaggregation (Iran ↔ Saudi specifically)",
        "Exclude 2020 / 2022 to test signal in quieter regimes",
        "LLM scoring layer on top of GDELT events",
    ], Inches(7.0), Inches(2.3), Inches(6.1), Inches(3.5),
        font_size=15, colour=WHITE, bullet="→ ")

    add_rect(s, Inches(0.4), Inches(5.9), Inches(12.5), Inches(1.25),
             RGBColor(0x0D, 0x1A, 0x30))
    add_textbox(s,
                "The absence of a strong signal is itself a finding: "
                "GDELT geopolitical data appears to be priced in quickly. "
                "This is consistent with market efficiency — and worth stating clearly.",
                Inches(0.6), Inches(6.0), Inches(12.2), Inches(1.0),
                font_size=14, colour=RGBColor(0xCC, 0xDD, 0xEE))

    set_notes(s, """\
[~30 s — conclusion, garder du temps]
En résumé : on a construit un pipeline complet de données alternatives — téléchargement GDELT, \
feature engineering, validation temporelle stricte, deux modèles, deux benchmarks, backtest.

Les modèles ne battent pas le benchmark en accuracy, mais ils prédisent les deux classes. \
Le signal géopolitique GDELT est réel et visuellement cohérent — mais le marché semble \
l'intégrer plus vite que 3 jours.

Ce résultat négatif n'est pas un échec : c'est une conclusion. L'absence de signal tradeable \
est compatible avec l'efficience des marchés sur une donnée publique et gratuite.

Merci. On est disponibles pour les questions.
""")


# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    slide_title(prs)
    slide_question(prs)
    slide_data(prs)
    slide_features(prs)
    slide_method(prs)
    slide_results(prs)
    slide_finding(prs)
    slide_limitations(prs)
    slide_conclusion(prs)

    out = Path(__file__).parent / "notebooks" / "presentation.pptx"
    prs.save(out)
    print(f"Saved → {out}")


if __name__ == "__main__":
    main()
